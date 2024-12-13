import streamlit as st
import pandas as pd
from io import BytesIO
from openpyxl import Workbook
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# ページの基本設定
st.set_page_config(page_title="アンケートツール", layout="centered")

# Bootstrapを適用
st.markdown(
    """
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <style>
        body {
            background-color: #f8f9fa;
        }
        .container {
            max-width: 600px;
            margin-top: 30px;
        }
        .form-label {
            font-weight: bold;
        }
    </style>
    """,
    unsafe_allow_html=True
)

# タイトル
st.markdown(
    """
    <div class="container">
        <h1 class="text-center text-primary">アンケートツール</h1>
        <p class="text-center">以下の質問に回答してください。</p>
    </div>
    """,
    unsafe_allow_html=True
)

# 質問リスト
questions = [
    "具体的に我々は何を提供するのですか？",
    "つまり、顧客が得る成果は何ですか？",
    "具体的には顧客は何を体験するのですか？",
    "つまり、我々は何者ですか？",
    "あなたのXXXを1〜10で評価してください。"
]

# 回答を格納するリスト
responses = {}

# フォーム
st.markdown("<div class='container'>", unsafe_allow_html=True)
for i, question in enumerate(questions):
    if "XXX" in question:
        response = st.slider(question, min_value=1, max_value=10, key=f"q{i+1}")
    else:
        st.markdown(f"""
            <div class="mb-3">
                <label for="q{i}" class="form-label">{question}</label>
                <textarea class="form-control" id="q{i}" rows="3" placeholder="ここに回答を入力してください"></textarea>
            </div>
        """, unsafe_allow_html=True)
        response = st.text_area(f"回答 {i+1}", key=f"q{i+1}", placeholder="ここに回答を入力してください")
    responses[question] = response
st.markdown("</div>", unsafe_allow_html=True)

# 回答のまとめ表示
if st.button("回答をまとめる"):
    st.markdown("<div class='container'>", unsafe_allow_html=True)
    st.markdown("<h2 class='text-center text-success'>回答のまとめ</h2>", unsafe_allow_html=True)

    # 回答をデータフレームに変換して表示
    summary = pd.DataFrame(list(responses.items()), columns=["質問", "回答"])
    st.table(summary)

    # エクセルファイル作成
    output = BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "アンケート結果"

    # ヘッダー行
    ws.append(["質問", "回答"])

    # データ行
    for question, answer in responses.items():
        ws.append([question, answer])

    wb.save(output)
    output.seek(0)

    # エクセルダウンロードボタン
    st.download_button(
        label="エクセルファイルをダウンロード",
        data=output,
        file_name="アンケート結果.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    # 数値質問のグラフ
    if "XXX" in responses:
        st.markdown("<h3 class='text-center text-info'>グラフ表示</h3>", unsafe_allow_html=True)
        plt.figure(figsize=(6, 4))
        plt.bar(["XXX"], [responses["あなたのXXXを1〜10で評価してください。"]], color="skyblue")
        plt.ylim(0, 10)
        plt.title("XXXの評価")
        plt.xlabel("評価項目")
        plt.ylabel("スコア")
        st.pyplot(plt)

    # PPT作成
    ppt_output = BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "アンケート結果"

    for question, answer in responses.items():
        content = slide.shapes.add_textbox(Inches(1), Inches(1 + list(responses.keys()).index(question) * 0.5), Inches(8), Inches(0.5))
        text_frame = content.text_frame
        text_frame.text = f"{question}: {answer}"

    prs.save(ppt_output)
    ppt_output.seek(0)

    # PPTダウンロードボタン
    st.download_button(
        label="PPTファイルをダウンロード",
        data=ppt_output,
        file_name="アンケート結果.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.markdown("</div>", unsafe_allow_html=True)
